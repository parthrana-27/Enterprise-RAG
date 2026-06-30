import os
import logging
import traceback
from typing import List, Dict, Any
from sqlalchemy.orm import sessionmaker

from app.core.embeddings import get_embeddings
from app.models.models import Document, DocumentChunk, AuditLog

logger = logging.getLogger("pipeline")

def recursive_character_splitter(text: str, chunk_size: int = 800, chunk_overlap: int = 150) -> List[str]:
    """
    Splits text into chunks of target chunk_size while maintaining chunk_overlap.
    Splits along paragraph, sentence, and word boundaries.
    """
    if not text:
        return []
    
    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = ""

    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        
        # If adding paragraph is within limits
        if len(current_chunk) + len(paragraph) + 2 <= chunk_size:
            if current_chunk:
                current_chunk += "\n\n" + paragraph
            else:
                current_chunk = paragraph
        else:
            # Current chunk is full, save it
            if current_chunk:
                chunks.append(current_chunk)
            
            # If paragraph itself is too large, split it by lines/sentences
            if len(paragraph) > chunk_size:
                sentences = paragraph.split(". ")
                current_chunk = ""
                for sentence in sentences:
                    sentence = sentence.strip()
                    if not sentence:
                        continue
                    if len(current_chunk) + len(sentence) + 2 <= chunk_size:
                        if current_chunk:
                            current_chunk += ". " + sentence
                        else:
                            current_chunk = sentence
                    else:
                        if current_chunk:
                            chunks.append(current_chunk)
                        current_chunk = sentence
                
                # If still too large, force split by length
                if len(current_chunk) > chunk_size:
                    for i in range(0, len(current_chunk), chunk_size - chunk_overlap):
                        chunks.append(current_chunk[i:i + chunk_size])
                    current_chunk = ""
            else:
                # Set paragraph as the new current chunk, preserving overlap if possible
                overlap_text = current_chunk[-(chunk_overlap):] if len(current_chunk) > chunk_overlap else ""
                current_chunk = (overlap_text + "\n\n" + paragraph).strip() if overlap_text else paragraph

    if current_chunk:
        chunks.append(current_chunk)

    return chunks

def parse_pdf(file_path: str) -> List[Dict[str, Any]]:
    """Parses PDF page-by-page using PyPDF2."""
    chunks = []
    try:
        import PyPDF2
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text = page.extract_text() or ""
                page_chunks = recursive_character_splitter(text)
                for chunk_text in page_chunks:
                    chunks.append({
                        "content": chunk_text,
                        "metadata": {"page": page_num + 1}
                    })
    except Exception as e:
        logger.error(f"Error parsing PDF: {e}")
        raise e
    return chunks

def parse_docx(file_path: str) -> List[Dict[str, Any]]:
    """Parses Word DOCX files paragraph by paragraph."""
    chunks = []
    try:
        import docx
        doc = docx.Document(file_path)
        full_text = []
        for p in doc.paragraphs:
            if p.text.strip():
                full_text.append(p.text)
        
        text = "\n\n".join(full_text)
        doc_chunks = recursive_character_splitter(text)
        for i, chunk_text in enumerate(doc_chunks):
            chunks.append({
                "content": chunk_text,
                "metadata": {"section": "body", "index": i + 1}
            })
    except Exception as e:
        logger.error(f"Error parsing DOCX: {e}")
        raise e
    return chunks

def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
    """Parses PowerPoint PPTX files slide by slide."""
    chunks = []
    try:
        import pptx
        prs = pptx.Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            text = "\n".join(slide_text)
            slide_chunks = recursive_character_splitter(text)
            for chunk_text in slide_chunks:
                chunks.append({
                    "content": chunk_text,
                    "metadata": {"slide": slide_num + 1}
                })
    except Exception as e:
        logger.error(f"Error parsing PPTX: {e}")
        raise e
    return chunks

def parse_xlsx(file_path: str) -> List[Dict[str, Any]]:
    """Parses Excel XLSX files sheet by sheet."""
    chunks = []
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_rows = []
            for row in sheet.iter_rows(values_only=True):
                # Filter None and convert row to string representation
                row_str = " | ".join([str(val) for val in row if val is not None])
                if row_str.strip():
                    sheet_rows.append(row_str)
            
            text = "\n".join(sheet_rows)
            sheet_chunks = recursive_character_splitter(text)
            for chunk_text in sheet_chunks:
                chunks.append({
                    "content": chunk_text,
                    "metadata": {"sheet": sheet_name}
                })
    except Exception as e:
        logger.error(f"Error parsing XLSX: {e}")
        raise e
    return chunks

def parse_text_or_markdown(file_path: str) -> List[Dict[str, Any]]:
    """Parses standard markdown or plaintext files."""
    chunks = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        
        doc_chunks = recursive_character_splitter(text)
        for i, chunk_text in enumerate(doc_chunks):
            chunks.append({
                "content": chunk_text,
                "metadata": {"index": i + 1}
            })
    except Exception as e:
        logger.error(f"Error parsing text file: {e}")
        raise e
    return chunks

def extract_chunks(file_path: str, file_type: str) -> List[Dict[str, Any]]:
    """Dispatches parsing based on file extensions."""
    file_type = file_type.lower().strip(".")
    if file_type == "pdf":
        return parse_pdf(file_path)
    elif file_type in ["docx", "doc"]:
        return parse_docx(file_path)
    elif file_type in ["pptx", "ppt"]:
        return parse_pptx(file_path)
    elif file_type in ["xlsx", "xls"]:
        return parse_xlsx(file_path)
    else:
        # Default to markdown / text parsing
        return parse_text_or_markdown(file_path)

def ingest_document_in_background(db_session_maker: sessionmaker, document_id: int, file_path: str):
    """
    Background worker function running in a separate thread/task.
    Loads file, splits it, embeds chunks, and saves to Postgres database.
    """
    db = db_session_maker()
    doc = db.query(Document).filter(Document.id == document_id).first()
    if not doc:
        logger.error(f"Ingestion failed: Document ID {document_id} not found in database.")
        db.close()
        return

    try:
        # 1. Update status to processing
        doc.status = "processing"
        db.commit()

        # 2. Extract text chunks and citation metadata
        logger.info(f"Extracting chunks from: {file_path} (Type: {doc.file_type})")
        parsed_chunks = extract_chunks(file_path, doc.file_type)

        if not parsed_chunks:
            raise ValueError("No extractable text found in this document.")

        # 3. Generate embeddings in batches for efficiency
        logger.info(f"Generating embeddings for {len(parsed_chunks)} chunks.")
        texts = [chunk["content"] for chunk in parsed_chunks]
        embeddings = get_embeddings(texts)

        # 4. Save chunks and embeddings to database
        db_chunks = []
        for index, chunk in enumerate(parsed_chunks):
            # Combine content and metadata info into chunk string for retrieval index
            meta = chunk["metadata"]
            
            db_chunk = DocumentChunk(
                document_id=doc.id,
                chunk_index=index,
                content=chunk["content"],
                embedding=embeddings[index],
                version=doc.version
            )
            db_chunks.append(db_chunk)

        db.add_all(db_chunks)
        
        # 5. Update document status
        doc.status = "completed"
        db.commit()

        # 6. Log success audit
        audit_log = AuditLog(
            user_id=doc.author_id,
            action="INGEST_SUCCESS",
            target_type="DOCUMENT",
            target_id=str(doc.id),
            details=f"Successfully ingested {doc.filename} (Version {doc.version}, {len(parsed_chunks)} chunks)"
        )
        db.add(audit_log)
        db.commit()
        logger.info(f"Document {document_id} ingestion completed successfully!")

    except Exception as e:
        logger.error(f"Document ingestion failed: {e}\n{traceback.format_exc()}")
        db.rollback()
        
        doc.status = "failed"
        doc.error_message = str(e)
        db.commit()

        # Log failure audit
        audit_log = AuditLog(
            user_id=doc.author_id,
            action="INGEST_FAILURE",
            target_type="DOCUMENT",
            target_id=str(doc.id),
            details=f"Failed to ingest {doc.filename}. Error: {str(e)}"
        )
        db.add(audit_log)
        db.commit()

    finally:
        db.close()
